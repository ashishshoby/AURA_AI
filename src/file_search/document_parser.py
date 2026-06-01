import os
import PyPDF2
import docx
import pptx
from bs4 import BeautifulSoup

def extract_text(file_path):
    """
    Extracts text from a file based on its extension.
    Supported extensions: .pdf, .docx, .pptx, .txt, .html
    """
    _, extension = os.path.splitext(file_path)
    extension = extension.lower()

    try:
        if extension == '.pdf':
            return extract_text_from_pdf(file_path)
        elif extension == '.docx':
            return extract_text_from_docx(file_path)
        elif extension == '.pptx':
            return extract_text_from_pptx(file_path)
        elif extension == '.txt':
            return extract_text_from_txt(file_path)
        elif extension == '.html':
            return extract_text_from_html(file_path)
        else:
            return None
    except Exception as e:
        print(f"Error extracting text from {file_path}: {e}")
        return None

def extract_text_from_pdf(file_path):
    """Extracts text from a PDF file."""
    text = ""
    with open(file_path, 'rb') as f:
        reader = PyPDF2.PdfReader(f)
        for page in reader.pages:
            text += page.extract_text() or ""
    return text

def extract_text_from_docx(file_path):
    """Extracts text from a DOCX file."""
    doc = docx.Document(file_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text

def extract_text_from_pptx(file_path):
    """Extracts text from a PPTX file."""
    prs = pptx.Presentation(file_path)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_txt(file_path):
    """Extracts text from a TXT file."""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        return f.read()

def extract_text_from_html(file_path):
    """Extracts text from an HTML file."""
    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
        soup = BeautifulSoup(f, 'html.parser')
        return soup.get_text()
